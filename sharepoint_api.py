from typing import List, Optional, Union
import pandas as pd
import requests
import os
import urllib.parse
import io
import platform
from io import BytesIO
from langchain_core.document_loaders import Blob
from langchain_core.documents.base import Document
from langchain_community.document_loaders.parsers.pdf import PyPDFParser
from langchain_core.document_loaders.base import BaseLoader
from docx import Document as DocxDocument
from pptx import Presentation


def ensure_directory_exists(file_path):
    directory = os.path.dirname(file_path)
    if not os.path.exists(directory):
        os.makedirs(directory, exist_ok=True)

def get_long_path(path):
    # Check if the operating system is Windows
    if platform.system() == 'Windows':
        # Apply the \\?\ prefix correctly to handle long paths on Windows
        return '\\\\?\\' + os.path.abspath(path).strip()
    else:
        # Return the normal path for Unix-based systems
        return os.path.abspath(path)

class SharePointClient:
    def __init__(self, tenant_id, client_id, client_secret, resource_url):
        self.tenant_id = tenant_id
        self.client_id = client_id
        self.client_secret = client_secret
        self.resource_url = resource_url
        self.base_url = f"https://login.microsoftonline.com/{tenant_id}/oauth2/v2.0/token"
        self.headers = {'Content-Type': 'application/x-www-form-urlencoded'}
        self.access_token = self.get_access_token()  # Initialize and store the access token upon instantiation

    def get_access_token(self):
        """
    This function retrieves an access token from Microsoft's OAuth2 endpoint.
    
    The access token is used to authenticate and authorize the application for 
    accessing Microsoft Graph API resources.

    Returns:
    str: The access token as a string. This token is used for authentication in subsequent API requests.
    """
        # Body for the access token request
        body = {
            'grant_type': 'client_credentials',
            'client_id': self.client_id,
            'client_secret': self.client_secret,
            'scope': self.resource_url + '.default'
        }
        response = requests.post(self.base_url, headers=self.headers, data=body)
        return response.json().get('access_token')  # Extract access token from the response

    def get_site_id(self, site_url):
        """
    This function retrieves the ID of a SharePoint site using the Microsoft Graph API.
    
    Parameters:
    site_url (str): The URL of the SharePoint site.

    Returns:
    str: The ID of the SharePoint site.
    """
        # Build URL to request site ID
        full_url = f'https://graph.microsoft.com/v1.0/sites/{site_url}'
        response = requests.get(full_url, headers={'Authorization': f'Bearer {self.access_token}'})
        return response.json().get('id')  # Return the site ID

    def get_drive_id(self, site_id):
        """
    This function retrieves the IDs and names of all drives associated with a specified SharePoint site.
    
    Parameters:
    site_id (str): The ID of the SharePoint site.

    Returns:
    list: A list of dictionaries. Each dictionary represents a drive on the SharePoint site.
          Each dictionary contains the following keys:
          - 'id': The ID of the drive.
          - 'name': The name of the drive.
    """
        # Retrieve drive IDs and names associated with a site
        drives_url = f'https://graph.microsoft.com/v1.0/sites/{site_id}/drives'
        response = requests.get(drives_url, headers={'Authorization': f'Bearer {self.access_token}'})
        drives = response.json().get('value', [])
        return [({'id': drive['id'], 'name': drive['name']}) for drive in drives]


    def get_folder_id(self, site_id, drive_id, folder_path):
        """
        This function retrieves the ID of a specified subfolder.
        
        Parameters:
        site_id (str): The ID of the site where the subfolder is located.
        drive_id (str): The ID of the drive where the subfolder is located.
        folder_path (str): The path of the subfolder whose ID is to be retrieved.

        Returns:
        str: The ID of the specified subfolder.
        """

        # Split the folder path into individual folders
        folders = folder_path.split('/')

        # Start with the root folder
        current_folder_id = 'root'

        # Loop through each folder in the path
        for folder_name in folders:
            # Build the URL to access the contents of the current folder
            folder_url = f'https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}/items/{current_folder_id}/children'
            response = requests.get(folder_url, headers={'Authorization': f'Bearer {self.access_token}'})
            items_data = response.json()

            # Loop through the items and find the folder
            for item in items_data['value']:
                if 'folder' in item and item['name'] == folder_name:
                    # Update the current folder ID and break the loop
                    current_folder_id = item['id']
                    break
            else:
                # If the folder was not found, return None
                return None

        # Return the ID of the final folder in the path
        return current_folder_id

    def list_folder_contents(self, site_id, drive_id, folder_id='root'):
        """
    This function lists the contents of a specific folder in a drive on a site.

    Parameters:
    site_id (str): The ID of the site.
    drive_id (str): The ID of the drive.
    folder_id (str, optional): The ID of the folder. Defaults to 'root'.

    Returns:
    list: A list of dictionaries. Each dictionary contains details about an item in the folder.
          The details include 'id', 'name', 'type', 'mimeType', 'uri', 'path', 'fullpath', 'filename', and 'url'.
    """
        items_list = []
        folder_contents_url = f'https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}/items/{folder_id}/children'
        while folder_contents_url:
                contents_response = requests.get(folder_contents_url, headers={'Authorization': f'Bearer {self.access_token}'})
                folder_contents = contents_response.json()
                for item in folder_contents.get('value', []):
                    path_parts = item['parentReference']['path'].split('root:')
                    path = path_parts[1] if len(path_parts) > 1 else ''
                    full_path = f"{path}/{item['name']}" if path else item['name']
                    
                    # Modifiez le site_web_url pour pointer vers l'élément spécifique
                    item_url = f'https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}/items/{item["id"]}'
                    response = requests.get(item_url, headers={'Authorization': f'Bearer {self.access_token}'})
                    item_data = response.json()
                    item_web_url = item_data.get('webUrl', '')

                    items_list.append({
                        'id': item['id'],
                        'name': item['name'],
                        'type': 'folder' if 'folder' in item else 'file',
                        'mimeType': item['file']['mimeType'] if 'file' in item else '',
                        'uri': item.get('@microsoft.graph.downloadUrl', ''),
                        'path': path,
                        'fullpath': full_path,
                        'filename': item['name'],
                        'url': item_web_url
                    })
                folder_contents_url = folder_contents.get('@odata.nextLink')
        return items_list

    
    def download_file(self, download_url, local_path, file_name):
        """
    This function downloads a file from a specified URL and saves it to a local path.
    
    Parameters:
    download_url (str): The URL of the file to be downloaded.
    local_path (str): The local path where the file will be saved.
    file_name (str): The name of the file to be saved.

    Returns:
    None. The function prints a success message if the file is downloaded and saved successfully, 
    or an error message if the download fails.
    """
        headers = {'Authorization': f'Bearer {self.access_token}'}
        response = requests.get(download_url, headers=headers)
        if response.status_code == 200:
            full_path = os.path.join(local_path, file_name)
            full_path = get_long_path(full_path)  # Apply the long path fix conditionally based on the OS
            ensure_directory_exists(full_path) 
            with open(full_path, 'wb') as file:
                file.write(response.content)
            # print(f"File downloaded: {full_path}")
        else:
            print(f"Failed to download {file_name}: {response.status_code} - {response.reason}")


    def download_folder_contents(self, site_id, drive_id, folder_id, local_folder_path, level=0):
        """
    This function recursively downloads all contents from a specified folder on a SharePoint site.
    
    Parameters:
    site_id (str): The ID of the SharePoint site.
    drive_id (str): The ID of the drive on the SharePoint site.
    folder_id (str): The ID of the folder whose contents are to be downloaded.
    local_folder_path (str): The local path where the downloaded files will be saved.
    level (int, optional): The current level of recursion (folder depth). Defaults to 0 for the root folder.

    Returns:
    None. The function saves the downloaded files to the specified local path and prints a success message for each downloaded file.
    """
        # Recursively download all contents from a folder
        folder_contents_url = f'https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}/items/{folder_id}/children'
        contents_headers = {'Authorization': f'Bearer {self.access_token}'}
        contents_response = requests.get(folder_contents_url, headers=contents_headers)
        folder_contents = contents_response.json()

        if 'value' in folder_contents:
            for item in folder_contents['value']:
                if 'folder' in item:
                    new_path = os.path.join(local_folder_path, item['name'])
                    if not os.path.exists(new_path):
                        os.makedirs(new_path)
                    self.download_folder_contents(site_id, drive_id, item['id'], new_path,
                                                  level + 1)  # Recursive call for subfolders
                elif 'file' in item:
                    file_name = item['name']
                    file_download_url = f"{self.resource_url}/v1.0/sites/{site_id}/drives/{drive_id}/items/{item['id']}/content"
                    self.download_file(file_download_url, local_folder_path, file_name)
    def download_file_contents(self, site_id, drive_id, file_id, local_save_path):
        """
        This function downloads the contents of a specified file from a SharePoint site and saves it to a local path.
        
        Parameters:
        site_id (str): The ID of the SharePoint site.
        drive_id (str): The ID of the drive on the SharePoint site.
        file_id (str): The ID of the file to be downloaded.
        local_save_path (str): The local path where the downloaded file will be saved.

        Returns:
        bool: True if the file was successfully downloaded and saved, False otherwise.
        """
        try:
            # Get the file details
            file_url = f'https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}/items/{file_id}'
            headers = {'Authorization': f'Bearer {self.access_token}'}
            response = requests.get(file_url, headers=headers)
            file_data = response.json()

            # Get the download URL and file name
            download_url = file_data['@microsoft.graph.downloadUrl']
            file_name = file_data['name']
            sharepoint_file_path = file_data['parentReference']['path']  # This is the SharePoint file path
            index = sharepoint_file_path.find(":/")

            # Extract everything after ":/"
            if index != -1:
                extracted_path = sharepoint_file_path[index+2:]  # Adding 2 to skip the characters ":/"
                local_save_path = local_save_path + "/" + extracted_path
                os.makedirs(local_save_path, exist_ok=True) # create loclal sub-folder
            else:
                extracted_path = ""
            # print(f"Downloading {file_name} from {extracted_path}")   
            
            # Download the file
            self.download_file(download_url, local_save_path, file_name)

            # If no exception was raised, the file download was successful
            return True

        except requests.exceptions.RequestException as e:
            print(f"Error downloading file: {file_name} err: {e}")
            return False

    def download_all_files(self, site_id, drive_id, local_folder_path, sharepoint_path="root"):
        """
        This method initiates the download of all files from a specific drive on a site.

        Args:
            site_id (str): The ID of the site from which files are to be downloaded.
            drive_id (str): The ID of the drive on the site from which files are to be downloaded.
            local_folder_path (str): The local path where the downloaded files should be stored.
        """
        try:
            if sharepoint_path != "root":
                folder_id = self.get_folder_id(site_id, drive_id, sharepoint_path)
            else:
                folder_id = sharepoint_path

            self.recursive_download(site_id, drive_id, folder_id, local_folder_path)
        except Exception as e:
            print(f"An error occurred while downloading files: {e}")

    def recursive_download(self, site_id, drive_id, folder_id, local_path):
        """
        This method downloads files from a folder and its subfolders recursively.

        Args:
            site_id (str): The ID of the site from which files are to be downloaded.
            drive_id (str): The ID of the drive on the site from which files are to be downloaded.
            folder_id (str): The ID of the folder from which files are to be downloaded.
            local_path (str): The local path where the downloaded files should be stored.
        """
        try:
            folder_contents = self.list_folder_contents(site_id, drive_id, folder_id)
            for item in folder_contents:
                sharepoint_path = item['path']
                sharepoint_path = sharepoint_path.lstrip('/')
                new_local_path = os.path.normpath(os.path.join(local_path, sharepoint_path))
                # Ensure the local directory exists before downloading
                #"data2\\BASIC STUDIES\\1. Cross Category
                #os.makedirs("data2\\BASIC STUDIES\\1. Cross Category", exist_ok=True)
                
                os.makedirs(new_local_path, exist_ok=True)
                if item['type'] == 'folder':
                    self.recursive_download(site_id, drive_id, item['id'], local_path)
                elif item['type'] == 'file':
                    # os.makedirs(os.path.dirname(new_local_path), exist_ok=True)
                    self.download_file(item['uri'], new_local_path, item['name'])
        except Exception as e:
            print(f"An error occurred while recursively downloading files:{new_local_path} {e}")


    def load_sharepoint_document(self, site_id, drive_id, file_id, file_name, file_type):
        """
    This function retrieves a document from a SharePoint site and loads it into memory using a custom loader based on the file type.
    
    Parameters:
    site_id (str): The ID of the SharePoint site.
    drive_id (str): The ID of the drive on the SharePoint site.
    file_id (str): The ID of the file to be loaded.
    file_name (str): The name of the file to be loaded.
    file_type (str): The MIME type of the file to be loaded.

    Returns:
    object: A custom loader object that can handle the content of the loaded file. The type of the loader depends on the file type.
    """
        # Get the download URL and the file name by querying the Microsoft Graph API
        file_url = f'https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}/items/{file_id}'
        headers = {'Authorization': f'Bearer {self.access_token}'}  # Use the stored access token for authorization
        response = requests.get(file_url, headers=headers)  # Make the HTTP request to get file details
        file_data = response.json()  # Parse the JSON response to get file data
        # print(file_data)

        download_url = file_data['@microsoft.graph.downloadUrl']  # Extract the direct download URL from the response

        # Get the file content from the download URL
        response = requests.get(download_url, headers=headers)  # Make the HTTP request to download the file

        # Create a BytesIO object from the response content, which allows for reading and writing bytes in memory
        stream = io.BytesIO(response.content)  # This is useful for handling binary data like files without saving to disk

        # Check the file type and use the appropriate custom loader to handle the file content
        if file_type == 'application/pdf':
            # Use CustomPDFLoader to handle PDF files; it initializes with the stream and file name
            loader = CustomPDFLoader(stream, file_name)
            return loader
        elif file_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            # Use CustomWordLoader for Word documents to handle and potentially split the document's content
            loader = CustomWordLoader(stream, file_name)
            return loader
        elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            # Use CustomPPTLoader for PowerPoint presentations to read and split the presentation into slides
            loader = CustomPPTLoader(stream, file_name)
            return loader
        elif file_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            # Use CustomExcelLoader for Excel spreadsheets to read and possibly split the sheets into separate parts
            loader = CustomExcelLoader(stream, file_name)
            return loader
        elif file_type in ['text/csv', 'text/plain']:
            # Use CustomTextLoader for plain text or CSV files to handle and split text as needed
            loader = CustomTextLoader(stream, file_name)
            return loader
        else:
            print(f"Unsupported file type: {file_type}")
            pass  # Placeholder for additional file types that may need to be implemented in the future



class CustomPDFLoader(BaseLoader):
    """
    This class is a custom loader for PDF files. It inherits from the BaseLoader class.
    
    The class is initialized with a binary stream of the PDF file, the file name, an optional password for protected PDFs, 
    and a flag indicating whether to extract images from the PDF. 
    
    The load method converts the binary stream into a Blob object, parses the PDF, and converts each page or segment 
    into a separate document object. The file name is added as metadata to each document.
    """
    def __init__(self, stream: BytesIO, filename: str, password: Optional[Union[str, bytes]] = None,
                 extract_images: bool = False):
        # Initialize with a binary stream, file name, optional password, and an image extraction flag
        self.stream = stream
        self.filename = filename
        # Initialize a PDF parser with optional password protection and image extraction settings
        self.parser = PyPDFParser(password=password, extract_images=extract_images)

    def load(self) -> List[Document]:
        # Convert the binary stream into a Blob object which is required by the parser
        blob = Blob.from_data(self.stream.getvalue())
        # Parse the PDF and convert each page or segment into a separate document object
        documents = list(self.parser.parse(blob))

        # Add the filename as metadata to each document for identification
        for doc in documents:
            doc.metadata.update({'source': self.filename})

        return documents

class CustomWordLoader(BaseLoader):
    """
    This class is a custom loader for Word documents. It extends the BaseLoader class and overrides its methods.
    It uses the python-docx library to parse Word documents and optionally splits the text into manageable documents.
    
    Attributes:
    stream (io.BytesIO): A binary stream of the Word document.
    filename (str): The name of the Word document.
    """
    def __init__(self, stream, filename: str):
        # Initialize with a binary stream and filename
        self.stream = stream
        self.filename = filename

    def load_and_split(self, text_splitter=None):
        # Use python-docx to parse the Word document from the binary stream
        doc = DocxDocument(self.stream)
        # Extract and concatenate all paragraph texts into a single string
        text = "\n".join([p.text for p in doc.paragraphs])

        # Check if a text splitter utility is provided
        if text_splitter is not None:
            # Use the provided splitter to divide the text into manageable documents
            split_text = text_splitter.create_documents([text])
        else:
            # Without a splitter, treat the entire text as one document
            split_text = [{'text': text, 'metadata': {'source': self.filename}}]

        # Add source metadata to each resulting document
        for doc in split_text:
            if isinstance(doc, dict):
                doc['metadata'] = {**doc.get('metadata', {}), 'source': self.filename}
            else:
                doc.metadata = {**doc.metadata, 'source': self.filename}

        return split_text

class CustomExcelLoader(BaseLoader):
    """
    This class is a custom loader for Excel files. It inherits from the BaseLoader class.
    
    The class takes a binary stream of an Excel file and a filename as input, and provides a method to load the Excel file into memory and split its content into separate documents based on the sheets in the workbook.
    
    Attributes:
    stream (io.BytesIO): A binary stream of the Excel file.
    filename (str): The name of the Excel file.
    """
    def __init__(self, stream, filename: str):
        # Initialize with a binary stream and filename
        self.stream = stream
        self.filename = filename

    def load_and_split(self, text_splitter=None):
        # Use pandas to load the Excel file from the binary stream
        xls = pd.ExcelFile(self.stream, engine='openpyxl')
        # Get the list of all sheet names in the workbook
        sheet_names = xls.sheet_names

        split_sheets = []
        for sheet in sheet_names:
            # Parse each sheet into a DataFrame
            df = xls.parse(sheet)
            # Convert the DataFrame to a single string with each cell value separated by new lines
            text = '\n'.join(df.values.astype(str).flatten().tolist())

            # Check if a text splitter is provided to further divide the sheet content
            if text_splitter is not None:
                # Use the splitter to create documents from the text
                split_text = text_splitter.create_documents([text])
                # Add metadata to each document
                for doc in split_text:
                    doc.metadata = {'source': self.filename, 'page': sheet}
                split_sheets.extend(split_text)
            else:
                # Without a splitter, treat the entire sheet text as one document
                doc = Document(text, metadata={'source': self.filename, 'page': sheet})
                split_sheets.append(doc)

        return split_sheets

class CustomPPTLoader(BaseLoader):
    """
    This class is a custom loader for PowerPoint files. It inherits from the BaseLoader class.
    
    The class takes a binary stream of a PowerPoint file and a filename as input, and provides a method to load the PowerPoint file into memory and split its content into separate documents based on the slides in the presentation.
    
    Attributes:
    stream (io.BytesIO): A binary stream of the PowerPoint file.
    filename (str): The name of the PowerPoint file.
    """
    def __init__(self, stream, filename):
        # Initialize with a binary stream and filename
        self.stream = stream
        self.filename = filename

    def load_and_split(self, text_splitter=None):
        # Use python-pptx to parse the PowerPoint file from the binary stream
        prs = Presentation(self.stream)
        # Prepare to collect documents
        documents = []

        # Iterate over each slide in the presentation
        for i, slide in enumerate(prs.slides):
            # Extract all text content from each slide
            slide_text = "\n".join([paragraph.text for shape in slide.shapes if shape.has_text_frame for paragraph in shape.text_frame.paragraphs])

            # Check if a text splitter is provided
            if text_splitter is None:
                # Treat each slide's text as a single document
                doc = {'text': slide_text, 'metadata': {'source': self.filename, 'page': i + 1}}
                # doc = {'text': slide_text, 'metadata': {'source': self.filename}}
                documents.append(doc)
            else:
                # Use the splitter to divide the slide text into smaller documents
                split_text = text_splitter.create_documents([slide_text])
                # Add metadata and collect each document
                for doc in split_text:
                    doc.metadata = {'source': self.filename, 'page': i + 1}
                    # doc.metadata = {'source': self.filename}
                documents.extend(split_text)

        return documents



import chardet

class CustomTextLoader(BaseLoader):
    """
    This class is a custom loader for text files. It inherits from the BaseLoader class.
    
    The class takes a binary stream of a text file and a filename as input, and provides a method to load the text file into memory and split its content into separate documents.
    
    Attributes:
    stream (io.BytesIO): A binary stream of the text file.
    filename (str): The name of the text file.
    """
    def __init__(self, stream, filename: str):
        self.stream = stream
        self.filename = filename

    def load_and_split(self, text_splitter=None):
        # Use chardet to detect the encoding of the stream
        rawdata = self.stream.read()
        result = chardet.detect(rawdata)
        text = rawdata.decode(result['encoding'])

        if text_splitter is not None:
            split_text = text_splitter.create_documents([text])
        else:
            split_text = [{'text': text, 'metadata': {'source': self.filename}}]

        for doc in split_text:
            if isinstance(doc, dict):
                doc['metadata'] = {**doc.get('metadata', {}), 'source': self.filename}
            else:
                doc.metadata = {**doc.metadata, 'source': self.filename}

        return split_text
