from typing import List, Optional, Union
import pandas as pd
import requests
import os
import io
from io import BytesIO
from langchain_core.document_loaders import Blob
from langchain_core.documents.base import Document
from langchain_community.document_loaders.parsers.pdf import PyPDFParser
from langchain_core.document_loaders.base import BaseLoader
from docx import Document as DocxDocument
from pptx import Presentation


class SharePointClient:
    def __init__(self, tenant_id, client_id, client_secret, resource_url):
        self.tenant_id = tenant_id
        self.client_id = client_id
        self.client_secret = client_secret
        self.resource_url = resource_url
        self.base_url = f"https://login.microsoftonline.com/{tenant_id}/oauth2/v2.0/token"
        self.headers = {'Content-Type': 'application/x-www-form-urlencoded'}
        self.access_token = self.get_access_token()  # Initialize and store the access token upon instantiation

    def get_access_token(self):
        # Body for the access token request
        body = {
            'grant_type': 'client_credentials',
            'client_id': self.client_id,
            'client_secret': self.client_secret,
            'scope': self.resource_url + '.default'
        }
        response = requests.post(self.base_url, headers=self.headers, data=body)
        return response.json().get('access_token')  # Extract access token from the response

    def get_site_id(self, site_url):
        # Build URL to request site ID
        full_url = f'https://graph.microsoft.com/v1.0/sites/{site_url}'
        response = requests.get(full_url, headers={'Authorization': f'Bearer {self.access_token}'})
        return response.json().get('id')  # Return the site ID

    def get_drive_id(self, site_id):
        # Retrieve drive IDs and names associated with a site
        drives_url = f'https://graph.microsoft.com/v1.0/sites/{site_id}/drives'
        response = requests.get(drives_url, headers={'Authorization': f'Bearer {self.access_token}'})
        drives = response.json().get('value', [])
        return [({'id': drive['id'], 'name': drive['name']}) for drive in drives]

    def get_folder_content(self, site_id, drive_id, folder_path='root'):
        # Get the contents of a folder
        folder_url = f'https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}/root/children'
        response = requests.get(folder_url, headers={'Authorization': f'Bearer {self.access_token}'})
        items_data = response.json()
        rootdir = []
        if 'value' in items_data:
            for item in items_data['value']:
                rootdir.append(({'id': item['id'], 'name': item['name']}))
        return rootdir

    # Recursive function to browse folders
    def list_folder_contents(self, site_id, drive_id, folder_id, level=0):
        # Get the contents of a specific folder
        folder_contents_url = f'https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}/items/{folder_id}/children'
        contents_headers = {'Authorization': f'Bearer {self.access_token}'}
        contents_response = requests.get(folder_contents_url, headers=contents_headers)
        folder_contents = contents_response.json()

        items_list = []  # List to store information

        if 'value' in folder_contents:
            for item in folder_contents['value']:
                # Split the path on 'root:' and take the second part, if it exists
                path_parts = item['parentReference']['path'].split('root:')
                path = path_parts[1] if len(path_parts) > 1 else ''
                full_path = f"{path}/{item['name']}" if path else item['name']

                if 'folder' in item:
                    # Add folder to list
                    items_list.append(
                        {'id': item['id'], 'name': item['name'], 'type': 'Folder', 'mimeType': None, 'uri': None,
                         'path': full_path})
                    # Recursive call for subfolders
                    items_list.extend(self.list_folder_contents(site_id, drive_id, item['id'], level + 1))
                elif 'file' in item:
                    # Add file to the list with its mimeType and uri
                    items_list.append(
                        {'id': item['id'], 'name': item['name'], 'type': 'File', 'mimeType': item['file']['mimeType'],
                         'uri': item['@microsoft.graph.downloadUrl'], 'path': full_path})

        return items_list

    def download_file(self, download_url, local_path, file_name):
        headers = {'Authorization': f'Bearer {self.access_token}'}
        response = requests.get(download_url, headers=headers)
        if response.status_code == 200:
            full_path = os.path.join(local_path, file_name)
            with open(full_path, 'wb') as file:
                file.write(response.content)
            print(f"File downloaded: {full_path}")
        else:
            print(f"Failed to download {file_name}: {response.status_code} - {response.reason}")

    def download_folder_contents(self, site_id, drive_id, folder_id, local_folder_path, level=0):
        # Recursively download all contents from a folder
        folder_contents_url = f'https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}/items/{folder_id}/children'
        contents_headers = {'Authorization': f'Bearer {self.access_token}'}
        contents_response = requests.get(folder_contents_url, headers=contents_headers)
        folder_contents = contents_response.json()

        if 'value' in folder_contents:
            for item in folder_contents['value']:
                if 'folder' in item:
                    new_path = os.path.join(local_folder_path, item['name'])
                    if not os.path.exists(new_path):
                        os.makedirs(new_path)
                    self.download_folder_contents(site_id, drive_id, item['id'], new_path,
                                                  level + 1)  # Recursive call for subfolders
                elif 'file' in item:
                    file_name = item['name']
                    file_download_url = f"{self.resource_url}/v1.0/sites/{site_id}/drives/{drive_id}/items/{item['id']}/content"
                    self.download_file(file_download_url, local_folder_path, file_name)

    def download_file_contents(self, site_id, drive_id, file_id, local_save_path):
        # Get the file details
        file_url = f'https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}/items/{file_id}'
        headers = {'Authorization': f'Bearer {self.access_token}'}
        response = requests.get(file_url, headers=headers)
        file_data = response.json()

        # Get the download URL and file name
        download_url = file_data['@microsoft.graph.downloadUrl']
        file_name = file_data['name']

        # Download the file
        self.download_file(download_url, local_save_path, file_name)

    def load_sharepoint_document(self, site_id, drive_id, file_id, file_name, file_type):
        # Obtenez l'URL de téléchargement et le nom du fichier
        file_url = f'https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}/items/{file_id}'
        headers = {'Authorization': f'Bearer {self.access_token}'}
        response = requests.get(file_url, headers=headers)
        file_data = response.json()
        download_url = file_data['@microsoft.graph.downloadUrl']

        # Obtenez la réponse
        response = requests.get(download_url, headers=headers)

        # Créez un objet BytesIO à partir de la réponse
        stream = io.BytesIO(response.content)

        if file_type == 'application/pdf':
            # Utilisez CustomPDFLoader pour charger et diviser le contenu du PDF
            loader = CustomPDFLoader(stream, file_name)
            return loader
        elif file_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            # Utilisez CustomWordLoader pour charger et diviser le contenu du document Word
            loader = CustomWordLoader(stream, file_name)
            return loader
        elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            # Utilisez CustomWordLoader pour charger et diviser le contenu du document Word
            loader = CustomPPTLoader(stream, file_name)
            return loader
        elif file_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            # Utilisez CustomWordLoader pour charger et diviser le contenu du document Word
            loader = CustomExcelLoader(stream, file_name)
            return loader
        elif file_type in ['text/csv', 'text/plain']:
            loader = CustomTextLoader(stream, file_name)
            return loader
        else:
            pass  # À implémenter


class CustomPDFLoader(BaseLoader):
    def __init__(self, stream: BytesIO, filename: str, password: Optional[Union[str, bytes]] = None,
                 extract_images: bool = False):
        self.stream = stream
        self.filename = filename
        self.parser = PyPDFParser(password=password, extract_images=extract_images)

    def load(self) -> List[Document]:
        blob = Blob.from_data(self.stream.getvalue())
        documents = list(self.parser.parse(blob))

        # Ajoutez les métadonnées à chaque document
        for doc in documents:
            doc.metadata.update({'source': self.filename})

        return documents


class CustomWordLoader(BaseLoader):
    def __init__(self, stream, filename: str):
        self.stream = stream
        self.filename = filename

    def load_and_split(self, text_splitter=None):
        # Utilisez python-docx pour lire le contenu du document Word
        doc = DocxDocument(self.stream)

        # Convertissez le contenu du document en texte
        text = "\n".join([p.text for p in doc.paragraphs])

        # Vérifiez si text_splitter est None
        if text_splitter is not None:
            # Si text_splitter est fourni, appliquez-le au texte
            split_text = text_splitter.create_documents([text])
        else:
            # Si text_splitter n'est pas fourni, utilisez simplement le texte tel quel
            split_text = [{'text': text, 'metadata': {'source': self.filename}}]

        # Ajoutez les métadonnées à chaque document
        for doc in split_text:
            if isinstance(doc, dict):
                doc['metadata'] = {**doc.get('metadata', {}), 'source': self.filename}
            else:
                doc.metadata = {**doc.metadata, 'source': self.filename}

        return split_text


class CustomExcelLoader(BaseLoader):
    def __init__(self, stream, filename: str):
        self.stream = stream
        self.filename = filename

    def load_and_split(self, text_splitter=None):
        # Utilisez pandas pour lire le contenu du fichier Excel
        xls = pd.ExcelFile(self.stream)
        sheet_names = xls.sheet_names

        split_sheets = []
        for sheet in sheet_names:
            df = xls.parse(sheet)

            # Convertissez le DataFrame en texte, chaque ligne du DataFrame est une ligne de la chaîne
            text = '\n'.join(df.values.astype(str).flatten().tolist())

            # Appliquez le text_splitter au texte si il est fourni
            if text_splitter is not None:
                split_text = text_splitter.create_documents([text])
                # Ajoutez les métadonnées à chaque document et ajoutez chaque document à split_sheets
                for doc in split_text:
                    doc.metadata = {'source': self.filename, 'page': sheet}
                split_sheets.extend(split_text)
            else:
                # Créez un objet Document avec le texte et les métadonnées
                doc = Document(text, metadata={'source': self.filename, 'page': sheet})
                split_sheets.append(doc)

        return split_sheets


class CustomPPTLoader(BaseLoader):
    def __init__(self, stream, filename):
        self.stream = stream
        self.filename = filename

    def load_and_split(self, text_splitter=None):
        # Utilisez python-pptx pour lire le contenu du fichier PowerPoint
        prs = Presentation(self.stream)

        # Créez une liste pour stocker les documents
        documents = []

        # Parcourez chaque diapositive
        for i, slide in enumerate(prs.slides):
            # Extraire le texte de la diapositive
            slide_text = "\n".join([paragraph.text for shape in slide.shapes if shape.has_text_frame for paragraph in
                                    shape.text_frame.paragraphs])

            # Vérifiez si text_splitter est None
            if text_splitter is None:
                # Créez un document pour chaque diapositive
                doc = {'text': slide_text, 'metadata': {'source': self.filename, 'page': i + 1}}
                documents.append(doc)
            else:
                # Appliquez le text_splitter au texte de la diapositive
                split_text = text_splitter.create_documents([slide_text])

                # Ajoutez les métadonnées à chaque document dans split_text
                for doc in split_text:
                    doc.metadata = {'source': self.filename, 'page': i + 1}

                # Ajoutez les documents à la liste
                documents.extend(split_text)

        return documents


class CustomTextLoader(BaseLoader):
    def __init__(self, stream, filename: str):
        self.stream = stream
        self.filename = filename

    def load_and_split(self, text_splitter=None):
        # Utilisez la bibliothèque standard de Python pour lire le contenu du fichier texte
        text = self.stream.read().decode('utf-8')

        # Vérifiez si text_splitter est None
        if text_splitter is not None:
            # Si text_splitter est fourni, appliquez-le au texte
            split_text = text_splitter.create_documents([text])
        else:
            # Si text_splitter n'est pas fourni, utilisez simplement le texte tel quel
            split_text = [{'text': text, 'metadata': {'source': self.filename}}]

        # Ajoutez les métadonnées à chaque document
        for doc in split_text:
            if isinstance(doc, dict):
                doc['metadata'] = {**doc.get('metadata', {}), 'source': self.filename}
            else:
                doc.metadata = {**doc.metadata, 'source': self.filename}

        return split_text
